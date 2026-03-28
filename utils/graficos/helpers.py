import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
import os
import matplotlib.font_manager as fm
from matplotlib.ticker import FuncFormatter 
from pptx import Presentation
from pathlib import Path

# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════


def _register_fonts(fonts_dir: str = None):
    """
    Registra todas las fuentes TTF/OTF de la carpeta fonts_dir
    en el font manager de matplotlib.
    Si fonts_dir es None usa data/raw/fonts relativo a este archivo.
    """
    if fonts_dir is None:
        fonts_dir = Path(__file__).parent.parent.parent / "data" / "raw" / "fonts"
    else:
        fonts_dir = Path(fonts_dir)

    if not fonts_dir.exists():
        print(f"⚠️  Carpeta de fuentes no encontrada: {fonts_dir}")
        return

    fuentes_encontradas = list(fonts_dir.glob("*.ttf")) + list(fonts_dir.glob("*.otf"))

    if not fuentes_encontradas:
        print(f"⚠️  No se encontraron fuentes TTF/OTF en {fonts_dir}")
        return

    for font_path in fuentes_encontradas:
        fm.fontManager.addfont(str(font_path))
        print(f"  ✓ Fuente registrada: {font_path.name}")

    # Limpiar caché para que matplotlib reconozca las nuevas fuentes
    fm._load_fontmanager(try_read_cache=False)
    print(f"Font manager recargado — {len(fuentes_encontradas)} fuentes registradas.")


def evolucion_inyeccion_bess(df_gx_real: pd.DataFrame, df_gx_real_comparacion: pd.DataFrame) -> pd.DataFrame:

    def _calcular_trimestral(df: pd.DataFrame) -> pd.DataFrame:
        df = df.copy()

        mask = (
            df["tipo"].astype(str).str.strip().str.lower().eq("bess") &
            df["subtipo"].astype(str).str.strip().str.contains("inye", case=False, na=False)
        )
        df = df[mask].copy()

        if df.empty:
            return pd.DataFrame(columns=["trimestre", "inyeccion_retiro", "anio"])

        df["fecha_hora"] = pd.to_datetime(df["fecha_hora"], errors="coerce")
        df["trimestre"]  = "Q" + df["fecha_hora"].dt.quarter.astype(str)   # ← "Q1".."Q4"
        df["anio"]       = df["fecha_hora"].dt.year.astype(int)             # ← int, no str

        return (
            df.groupby(["trimestre", "anio"], as_index=False)["inyeccion_retiro"]
            .sum()
        )

    df_estudio     = _calcular_trimestral(df_gx_real)
    df_comparacion = _calcular_trimestral(df_gx_real_comparacion)

    # Debug — eliminar en producción
    print("Debug evolucion_inyeccion_bess")
    print("df_estudio:\n",     df_estudio.head())
    print("df_comparacion:\n", df_comparacion.head())

    resultado = pd.concat([df_estudio, df_comparacion], ignore_index=True)
    resultado["inyeccion_retiro"] = resultado["inyeccion_retiro"].astype(float)
    return resultado


def evolucion_vertimiento(df_vertimientos, df_vertimientos_comparacion):

    def _calcular_trimestral(df: pd.DataFrame, tag="") -> pd.DataFrame:
        df = df.copy()

        if df.empty:
            return pd.DataFrame(columns=["trimestre", "vertimiento", "anio"])

        # Usar columnas anio/mes directas si están disponibles (vienen de hora_mensual)
        if "anio" in df.columns and "mes" in df.columns:
            df["anio"]      = df["anio"].astype(int)
            df["trimestre"] = "Q" + ((df["mes"].astype(int) - 1) // 3 + 1).astype(str)
        else:
            # Fallback: parsear desde periodo
            df["periodo"]   = pd.to_datetime(df["periodo"], errors="coerce")
            df["trimestre"] = "Q" + df["periodo"].dt.quarter.astype(str)
            df["anio"]      = df["periodo"].dt.year.astype(int)

        print(f"[{tag}] años únicos: {sorted(df['anio'].unique())}")
        print(f"[{tag}] trimestres únicos: {sorted(df['trimestre'].unique())}")

        return (
            df.groupby(["trimestre", "anio"], as_index=False)["vertimiento"]
            .sum()
        )
    df_estudio     = _calcular_trimestral(df_vertimientos,            "estudio")
    df_comparacion = _calcular_trimestral(df_vertimientos_comparacion, "comparacion")

    return pd.concat([df_estudio, df_comparacion], ignore_index=True)


def _setup_theme(font_dict, font_family_dict, grid_alph, grid_lw, edge_color, legend_alpha):
    """Aplica el tema global Seaborn coherente con el estilo día típico."""
    _register_fonts()
    plt.rcParams["font.family"] = font_family_dict
    sns.set_theme(
        style="whitegrid",
        context="notebook",
        rc={
            "font.family":        font_family_dict,
            "axes.titlesize":     12,
            "axes.titleweight":   "bold",
            "axes.labelsize":     10,
            "xtick.labelsize":    8,
            "ytick.labelsize":    8,
            "grid.alpha":         grid_alph,
            "grid.linewidth":     grid_lw,
            "axes.edgecolor":     edge_color,
            "axes.linewidth":     0.9,
            "figure.facecolor":   "none",
            "axes.facecolor":     "white",
            "legend.frameon":     True,
            "legend.framealpha":  legend_alpha,
            "text.color":         font_dict,
            "axes.labelcolor":    font_dict,
            "xtick.color":        font_dict,
            "ytick.color":        font_dict,
        },
    )


def _fmt_thousands(x, pos):
    try:
        return f"{x:,.0f}".replace(",", ".")
    except Exception:
        return str(x)


def _guardar_fig(fig, path, dpi=300):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".svg":
        fig.savefig(path, format="svg", bbox_inches="tight")
    else:
        fig.savefig(
            path,
            dpi=dpi,
            bbox_inches="tight",
            pad_inches=0.06,
            transparent=True,
        )
    plt.close(fig)


def _estilo_leyenda(leg, font_dict, font_family_dict, edge_color, legend_alpha):
    leg.get_frame().set_alpha(legend_alpha)
    leg.get_frame().set_edgecolor(edge_color)
    for text in leg.get_texts():
        text.set_color(font_dict)
        text.set_fontfamily(font_family_dict)
    if leg.get_title():
        leg.get_title().set_color(font_dict)
        leg.get_title().set_fontfamily(font_family_dict)


def _estilo_ax(ax, grid_alpha, grid_lw, font_dict):
    """Aplica estilo base consistente a un eje."""
    ax.grid(True,  axis="y", alpha=grid_alpha, linewidth=grid_lw, color="#CCCCCC")
    ax.grid(False, axis="x")
    ax.yaxis.set_major_formatter(FuncFormatter(_fmt_thousands))
    ax.tick_params(axis="x", labelsize=8, colors=font_dict)
    ax.tick_params(axis="y", labelsize=8, colors=font_dict)
    sns.despine(ax=ax, top=True, right=True)


def listar_shapes(pptx_path, slide_idx):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    for i, shape in enumerate(slide.shapes):
        tipo = "TABLA" if shape.has_table else shape.shape_type
        print(f"  [{i}] nombre='{shape.name}'  tipo={tipo}")


def render_table_image(df, title, out_path, font_dict, font_family_dict, muted_dict, edge_color,
                       figsize=(4.49, 3.68), font_scale=1.0, dpi=150, top=10):
    """Renderiza un DataFrame como imagen PNG con el estilo visual del reporte."""
    fs_cell = round(8 * font_scale)

    COLOR_HEADER_BG = muted_dict["c1"]
    COLOR_HEADER_FG = font_dict
    COLOR_ROW       = "#FFFFFF"
    COLOR_BORDER    = edge_color

    df_show = df.copy()
    for col in df_show.columns:
        if pd.api.types.is_numeric_dtype(df_show[col]):
            df_show[col] = df_show[col].map(
                lambda v: f"{v:,.0f}".replace(",", ".") if pd.notnull(v) else ""
            )
    if len(df_show) > top:
        df_show = df_show.head(top)

    nrows, ncols = df_show.shape

    fig, ax = plt.subplots(figsize=figsize)
    ax.axis("off")
    fig.patch.set_facecolor("none")

    table = ax.table(
        cellText=df_show.values,
        colLabels=df_show.columns.tolist(),
        loc="center",
        cellLoc="center",
        colLoc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(fs_cell)
    table.scale(1, 1.4)

    for (r, c), cell in table.get_celld().items():
        cell.set_linewidth(0.4)
        cell.set_edgecolor(COLOR_BORDER)
        if r == 0:
            cell.set_facecolor(COLOR_HEADER_BG)
            cell.set_text_props(weight="bold", color=COLOR_HEADER_FG,
                                fontfamily=font_family_dict, fontsize=fs_cell)
        else:
            cell.set_facecolor(COLOR_ROW)
            cell.set_text_props(color=font_dict, fontfamily=font_family_dict, fontsize=fs_cell)

    table.auto_set_column_width(col=list(range(ncols)))
    fig.subplots_adjust(left=0.02, right=0.98, top=0.98, bottom=0.02)
    _guardar_fig(fig, out_path, dpi=dpi)