import os
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from matplotlib.ticker import FuncFormatter
from matplotlib.lines import Line2D
from matplotlib.colors import LinearSegmentedColormap
import geopandas as gpd
import seaborn as sns
import numpy as np
from pptx import Presentation
import calendar
from utils.inserta_texto_ppt import insertar_top_vertimiento
from utils.evoluciones_bess import evolucion_inyeccion_bess
from utils.evoluciones_vertimientos import evolucion_vertimiento

# ══════════════════════════════════════════════════════════════════════════════
# CONSTANTES GLOBALES DE ESTILO
# ══════════════════════════════════════════════════════════════════════════════

FONT_FAMILY  = "Candara"
FONT_COLOR   = "#003366"
GRID_ALPHA   = 0.18
GRID_LW      = 0.7
EDGE_COLOR   = "#D0D5DD"
LEGEND_ALPHA = 0.9

# Paleta Seaborn Muted — orden fijo para consistencia entre gráficos
MUTED = {
    "c1": "#9EC8E8",   # Azul pizarra   — año reciente / serie principal
    "c2": "#F4B89A",   # Naranja tostado — año anterior / serie secundaria
    "c3": "#A8DDA5",   # Verde salvia
    "c4": "#E8A5A5",   # Rojo arcilla
    "c5": "#C4A8D4",   # Violeta suave
    "c6": "#C4A882",   # Café rosado
}

# Colores específicos por tecnología (día típico — paleta propia clara)
COLOR_TECNOLOGIA = {
    "Solar":          "#F6C48E",
    "Eólica":         "#A8D5BA",
    "Hidro":          "#8EC5FF",
    "Geotérmica":     "#D9C2A3",
    "Térmica":        "#C9C2E6",
    "BESS Inyección": "#D96C6C",
    "BESS Retiro":    "#6FA8DC",
}

SHP_REGIONES = r"C:\Users\pinzunza\OneDrive - Ministerio de Energia\Escritorio\escritorio desrodenado\capas\Comunas\COMUNAS_NACIONAL.shp"

BAR_POINTS = {
    "Barra crucero 200kV":       {"lon": -69.5677773900849, "lat": -22.27773471974709},
    "Barra Pan de Azucar 220kV": {"lon": -71.100,           "lat": -29.900},
    "Barra Quillota 220kV":      {"lon": -71.260,           "lat": -32.880},
    "Barra Alto Jahuel 500kV":   {"lon": -70.630,           "lat": -33.720},
    "Barra Puerto Montt 220kV":  {"lon": -72.940,           "lat": -41.470},
    "Barra Charrua 500kV":       {"lon": -72.940,           "lat": -41.470},
}


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _setup_theme():
    """Aplica el tema global Seaborn coherente con el estilo día típico."""
    plt.rcParams["font.family"] = FONT_FAMILY
    sns.set_theme(
        style="whitegrid",
        context="notebook",
        rc={
            "font.family":        FONT_FAMILY,
            "axes.titlesize":     12,
            "axes.titleweight":   "bold",
            "axes.labelsize":     10,
            "xtick.labelsize":    8,
            "ytick.labelsize":    8,
            "grid.alpha":         GRID_ALPHA,
            "grid.linewidth":     GRID_LW,
            "axes.edgecolor":     EDGE_COLOR,
            "axes.linewidth":     0.9,
            "figure.facecolor":   "none",
            "axes.facecolor":     "white",
            "legend.frameon":     True,
            "legend.framealpha":  LEGEND_ALPHA,
            "text.color":         FONT_COLOR,
            "axes.labelcolor":    FONT_COLOR,
            "xtick.color":        FONT_COLOR,
            "ytick.color":        FONT_COLOR,
        },
    )


def _fmt_thousands(x, pos):
    try:
        return f"{x:,.0f}".replace(",", ".")
    except Exception:
        return str(x)


def _guardar_fig(fig, path, dpi=300):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".svg":
        fig.savefig(path, format="svg", bbox_inches="tight")
    else:
        fig.savefig(
            path,
            dpi=dpi,
            bbox_inches="tight",
            pad_inches=0.06,
            transparent=True,   # ← fondo transparente
        )
    plt.close(fig)


def _estilo_leyenda(leg):
    leg.get_frame().set_alpha(LEGEND_ALPHA)
    leg.get_frame().set_edgecolor(EDGE_COLOR)
    for text in leg.get_texts():
        text.set_color(FONT_COLOR)
        text.set_fontfamily(FONT_FAMILY)
    if leg.get_title():
        leg.get_title().set_color(FONT_COLOR)
        leg.get_title().set_fontfamily(FONT_FAMILY)


def _estilo_ax(ax):
    """Aplica estilo base consistente a un eje."""
    ax.grid(True,  axis="y", alpha=GRID_ALPHA, linewidth=GRID_LW, color="#CCCCCC")
    ax.grid(False, axis="x")
    ax.yaxis.set_major_formatter(FuncFormatter(_fmt_thousands))
    ax.tick_params(axis="x", labelsize=8, colors=FONT_COLOR)
    ax.tick_params(axis="y", labelsize=8, colors=FONT_COLOR)
    sns.despine(ax=ax, top=True, right=True)


def listar_shapes(pptx_path, slide_idx):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    for i, shape in enumerate(slide.shapes):
        tipo = "TABLA" if shape.has_table else shape.shape_type
        print(f"  [{i}] nombre='{shape.name}'  tipo={tipo}")


# ══════════════════════════════════════════════════════════════════════════════
# MAPA DE REGIONES
# ══════════════════════════════════════════════════════════════════════════════

def generar_mapa_regiones(shp_path, target_crs="EPSG:32719"):
    CUT_COM_EXCLUIR  = [5201, 5104]
    REGIONES_INCLUIR = [15, 1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 14, 16, 13]

    gdf = gpd.read_file(shp_path)
    gdf["CUT_COM"] = gdf["CUT_COM"].astype(str).str.strip().astype(int)
    gdf = gdf.loc[~gdf["CUT_COM"].isin(CUT_COM_EXCLUIR)].copy()
    gdf["CUT_REG"] = gdf["CUT_REG"].astype(str).str.strip().astype(int)
    gdf = gdf.dissolve(by="CUT_REG", aggfunc="first").reset_index()
    if gdf.crs is None:
        gdf = gdf.set_crs("EPSG:4326", allow_override=True)
    gdf = gdf.to_crs(target_crs)
    gdf = gdf.loc[gdf["CUT_REG"].isin(REGIONES_INCLUIR)].copy()
    return gdf.sort_values("CUT_REG").reset_index(drop=True)


# ══════════════════════════════════════════════════════════════════════════════
# TABLA IMAGEN
# ══════════════════════════════════════════════════════════════════════════════

def render_table_image(df, title, out_path, figsize=(8.2, 4.0), dpi=300, top=10):
    COLOR_HEADER_BG = MUTED["c2"]
    COLOR_HEADER_FG = FONT_COLOR
    COLOR_ROW_ODD   = "#FFFFFF"
    COLOR_ROW_EVEN  = "#F0F4FA"
    COLOR_BORDER    = EDGE_COLOR

    df_show = df.copy()
    for col in df_show.columns:
        if pd.api.types.is_numeric_dtype(df_show[col]):
            df_show[col] = df_show[col].map(
                lambda v: f"{v:,.0f}".replace(",", ".") if pd.notnull(v) else ""
            )
    if len(df_show) > top:
        df_show = df_show.head(top)

    nrows, ncols = df_show.shape
    fig, ax = plt.subplots(figsize=figsize)
    ax.axis("off")
    fig.patch.set_facecolor("#FFFFFF")

    ax.set_title(title, fontsize=12, fontweight="bold", color=FONT_COLOR,
                 fontfamily=FONT_FAMILY, pad=3, loc="center")

    table = ax.table(cellText=df_show.values, colLabels=df_show.columns.tolist(),
                     loc="center", cellLoc="center", colLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1, 1.35)

    for (r, c), cell in table.get_celld().items():
        cell.set_linewidth(0.4)
        cell.set_edgecolor(COLOR_BORDER)
        if r == 0:
            cell.set_facecolor(COLOR_HEADER_BG)
            cell.set_text_props(weight="bold", color=COLOR_HEADER_FG,
                                fontfamily=FONT_FAMILY, fontsize=8)
        else:
            cell.set_facecolor(COLOR_ROW_ODD if r % 2 else COLOR_ROW_EVEN)
            cell.set_text_props(color=FONT_COLOR, fontfamily=FONT_FAMILY, fontsize=8)

    table.auto_set_column_width(col=list(range(ncols)))
    fig.subplots_adjust(left=0.02, right=0.98, top=0.88, bottom=0.02)
    _guardar_fig(fig, out_path, dpi=dpi)


# ══════════════════════════════════════════════════════════════════════════════
# 1) CMG CON MAPA
# ══════════════════════════════════════════════════════════════════════════════

def graficar_cmg_con_mapa(
    df_cmg,
    df_cmg_comparacion,
    gdf_reg,
    bar_points,
    out_path,
    figsize=(7.2, 3.2),
    dpi=300
):
    import calendar

    df_c = df_cmg.copy()
    df_c["fecha_hora"] = pd.to_datetime(df_c["fecha_hora"], errors="coerce")
    df_c = df_c.dropna(subset=["fecha_hora", "CMG_PESO_KWH", "nombre_cmg"]).copy()

    df_comp = df_cmg_comparacion.copy()
    df_comp["fecha_hora"] = pd.to_datetime(df_comp["fecha_hora"], errors="coerce")
    df_comp = df_comp.dropna(subset=["fecha_hora", "CMG_PESO_KWH", "nombre_cmg"]).copy()

    if df_c.empty and df_comp.empty:
        fig, ax = plt.subplots(figsize=figsize)  # ← aplicado
        ax.text(0.5, 0.5, "Sin datos para gráfico CMG", ha="center", va="center", fontsize=11)
        ax.axis("off")
        _guardar_fig(fig, out_path, dpi=dpi)
        return

    anio_estudio     = str(df_c["fecha_hora"].dt.year.mode().iloc[0])    if not df_c.empty    else "Año estudio"
    anio_comparacion = str(df_comp["fecha_hora"].dt.year.mode().iloc[0]) if not df_comp.empty else "Año anterior"

    for df_ in [df_c, df_comp]:
        if not df_.empty:
            df_["mes_num"] = df_["fecha_hora"].dt.month

    if not df_c.empty:
        df_c = (
            df_c.groupby(["mes_num", "nombre_cmg"], as_index=False)["CMG_PESO_KWH"]
            .mean()
            .sort_values(["nombre_cmg", "mes_num"])
        )

    if not df_comp.empty:
        df_comp = (
            df_comp.groupby(["mes_num", "nombre_cmg"], as_index=False)["CMG_PESO_KWH"]
            .mean()
            .sort_values(["nombre_cmg", "mes_num"])
        )

    barras  = sorted(
        set(df_c["nombre_cmg"].dropna().unique()).union(
            set(df_comp["nombre_cmg"].dropna().unique())
        )
    )
    palette   = sns.color_palette("pastel", n_colors=len(barras))
    color_map = dict(zip(barras, palette))

    fig, ax = plt.subplots(figsize=figsize)  # ← aplicado

    ax_map = fig.add_axes([1.03, 0.2, 0.16, 0.6])
    ax_map.grid(False)

    for barra in barras:
        g = df_c[df_c["nombre_cmg"] == barra].copy()
        if g.empty:
            continue
        ax.plot(g["mes_num"], g["CMG_PESO_KWH"],
                linestyle="-", linewidth=1.8, color=color_map[barra], alpha=0.95, zorder=3)

    for barra in barras:
        g = df_comp[df_comp["nombre_cmg"] == barra].copy()
        if g.empty:
            continue
        ax.plot(g["mes_num"], g["CMG_PESO_KWH"],
                linestyle="--", linewidth=1.5, color=color_map[barra], alpha=0.95, zorder=2)

    meses = np.arange(1, 13)
    ax.set_xticks(meses)
    ax.set_xticklabels([calendar.month_abbr[m] for m in meses], fontsize=8)
    ax.set_xlim(1, 12)
    ax.set_title("Costo marginal promedio por mes ($/kWh)", fontsize=11, fontweight="bold")
    ax.set_xlabel("")
    ax.set_ylabel("CMG ($/kWh)", fontsize=9)
    ax.tick_params(axis="y", labelsize=8)
    ax.grid(False)
    ax.yaxis.grid(True, alpha=0.18, linewidth=0.8)
    ax.xaxis.grid(False)
    sns.despine(ax=ax)

    gdf_reg_plot = gdf_reg.sort_values("CUT_REG").reset_index(drop=True)
    cmap = LinearSegmentedColormap.from_list("pastel_orange_blue", ["#F6B38E", "#8EC5FF"])
    colors = [cmap(i) for i in np.linspace(0, 1, len(gdf_reg_plot))]
    gdf_reg_plot["_color"] = colors
    gdf_reg_plot.plot(ax=ax_map, color=gdf_reg_plot["_color"],
                      edgecolor="white", linewidth=0.6, zorder=1)

    rows = [{"nombre_cmg": b, "lon": bar_points[b]["lon"], "lat": bar_points[b]["lat"]}
            for b in barras if b in bar_points]
    if rows:
        pts = gpd.GeoDataFrame(
            rows,
            geometry=gpd.points_from_xy([r["lon"] for r in rows], [r["lat"] for r in rows]),
            crs="EPSG:4326"
        ).to_crs(gdf_reg.crs)
        for _, r in pts.iterrows():
            ax_map.scatter(r.geometry.x, r.geometry.y, s=28, marker="o",
                           color=color_map[r["nombre_cmg"]], edgecolor="white",
                           linewidth=0.6, zorder=5)

    minx, miny, maxx, maxy = gdf_reg.total_bounds
    dx, dy = maxx - minx, maxy - miny
    ax_map.set_xlim(minx - 0.08*dx, maxx + 0.08*dx)
    ax_map.set_ylim(miny - 0.03*dy, maxy + 0.03*dy)
    ax_map.set_axis_off()
    ax_map.set_aspect("equal")
    ax_map.set_facecolor(ax.get_facecolor())
    ax_map.patch.set_alpha(1.0)
    for spine in ax_map.spines.values():
        spine.set_visible(True)
        spine.set_color("#D0D5DD")
        spine.set_linewidth(0.7)

    handles_barras = [
        Line2D([0], [0], color=color_map[b], linewidth=1.8, marker="o",
               markersize=4.5, markerfacecolor=color_map[b], markeredgecolor="white", label=b)
        for b in barras
    ]
    leg1 = ax.legend(handles=handles_barras, title="Barras CMG", loc="lower left",
                     fontsize=7, title_fontsize=8, frameon=True, ncol=2, borderaxespad=0.6)
    leg1.get_frame().set_alpha(0.9)
    ax.add_artist(leg1)

    handles_estilo = [
        Line2D([0], [0], color="#667085", linewidth=1.8, linestyle="-",  label=anio_estudio),
        Line2D([0], [0], color="#667085", linewidth=1.5, linestyle="--", label=anio_comparacion),
    ]
    leg2 = ax.legend(handles=handles_estilo, title="Período", loc="upper left",
                     fontsize=7, title_fontsize=8, frameon=True, ncol=1, borderaxespad=0.6)
    leg2.get_frame().set_alpha(0.9)

    fig.subplots_adjust(left=0.08, right=0.98, top=0.84, bottom=0.24)
    _guardar_fig(fig, out_path, dpi=dpi)


# ══════════════════════════════════════════════════════════════════════════════
# 2) DÍA TÍPICO
# ══════════════════════════════════════════════════════════════════════════════

def graficar_gx_tipico(df_dia_tipico, dia_tipico_comparacion,
                        fecha_tipica, fecha_tipica_comparacion,
                        out_path, figsize=(10.74, 5.02), dpi=300):

    def _preparar_df(df):
        if df.empty:
            return df
        df = df.copy()
        df["inyeccion_retiro"] = pd.to_numeric(df["inyeccion_retiro"], errors="coerce")
        df = df.dropna(subset=["inyeccion_retiro", "tipo"]).copy()

        if "hora_decimal" not in df.columns:
            if {"hora", "minuto"}.issubset(df.columns):
                df["hora"]   = pd.to_numeric(df["hora"],   errors="coerce").fillna(0)
                df["minuto"] = pd.to_numeric(df["minuto"], errors="coerce").fillna(0)
                df["hora_decimal"] = df["hora"] + df["minuto"] / 60.0
            else:
                df["fecha_hora"]   = pd.to_datetime(df["fecha_hora"], errors="coerce")
                df["hora_decimal"] = df["fecha_hora"].dt.hour + df["fecha_hora"].dt.minute / 60.0

        df["tipo"]    = df["tipo"].fillna("Sin clasificar").astype(str).str.strip()
        df["subtipo"] = df["subtipo"].fillna("-").astype(str).str.strip()

        rename_tipo = {
            "Eólicas": "Eólica", "Solar": "Solar", "Solares": "Solar",
            "Hidroeléctrica": "Hidro", "Hidroeléctricas": "Hidro", "Hidro": "Hidro",
            "Térmica": "Térmica", "Térmicas": "Térmica", "Termica": "Térmica", "Termicas": "Térmica",
            "Geotérmica": "Geotérmica", "Geotermia": "Geotérmica",
            "Bess": "BESS", "BESS": "BESS",
        }
        df["tipo_plot"]     = df["tipo"].replace(rename_tipo)
        df["categoria_plot"] = df["tipo_plot"]
        mask_bess   = df["tipo_plot"].eq("BESS")
        mask_retiro = df["subtipo"].str.contains("Retiro", case=False, na=False)
        mask_iny    = df["subtipo"].str.contains("Inye",   case=False, na=False)
        df.loc[mask_bess & mask_retiro, "categoria_plot"] = "BESS Retiro"
        df.loc[mask_bess & mask_iny,    "categoria_plot"] = "BESS Inyección"
        return df

    def _construir_pivot(df):
        if df.empty:
            return pd.DataFrame()
        df_plot = (df.groupby(["hora_decimal", "categoria_plot"], as_index=False)
                   ["inyeccion_retiro"].sum())
        orden = ["Solar", "Eólica", "Hidro", "Geotérmica", "Térmica", "BESS Inyección", "BESS Retiro"]
        pivot = (df_plot.pivot(index="hora_decimal", columns="categoria_plot",
                               values="inyeccion_retiro")
                 .fillna(0).sort_index())
        presentes  = [t for t in orden if t in pivot.columns]
        restantes  = sorted([t for t in pivot.columns if t not in presentes])
        return pivot[presentes + restantes]

    def _dibujar_areas(ax, pivot, fecha_label):
        if pivot.empty or len(pivot.index) <= 1 or np.isclose(pivot.to_numpy().sum(), 0):
            ax.text(0.5, 0.5, "Sin datos suficientes", ha="center", va="center",
                    fontsize=10, color=FONT_COLOR)
            ax.axis("off")
            return

        x        = pivot.index.values
        cols_pos = [c for c in pivot.columns if pivot[c].max() > 0]
        cols_neg = [c for c in pivot.columns if pivot[c].min() < 0]

        if cols_pos:
            ax.stackplot(x, [pivot[c].clip(lower=0).values for c in cols_pos],
                         labels=cols_pos,
                         colors=[COLOR_TECNOLOGIA.get(c, "#D9E2EC") for c in cols_pos],
                         alpha=0.95, linewidth=0.6)
        if cols_neg:
            ax.stackplot(x, [pivot[c].clip(upper=0).values for c in cols_neg],
                         labels=cols_neg,
                         colors=[COLOR_TECNOLOGIA.get(c, "#D9E2EC") for c in cols_neg],
                         alpha=0.95, linewidth=0.6)

        ax.axhline(0, color=EDGE_COLOR, linewidth=0.8)
        ax.text(0.01, 1.13, fecha_label, transform=ax.transAxes,
                fontsize=8, color="#667085", ha="left", va="bottom", fontstyle="italic",
                fontfamily=FONT_FAMILY)
        ax.set_title("Generación diaria típica por tecnología", fontsize=10,
                     fontweight="bold", pad=4, color=FONT_COLOR)
        ax.set_xlabel("Hora del día", fontsize=9, color=FONT_COLOR)
        ax.set_ylabel("Generación (mWh)", fontsize=9, color=FONT_COLOR)
        _estilo_ax(ax)
        xticks = np.arange(0, 25, 4)
        ax.set_xticks(xticks)
        ax.set_xlim(0, 24)
        ax.set_xticklabels([f"{int(h):02d}:00" for h in xticks])

    pivot_est  = _construir_pivot(_preparar_df(df_dia_tipico))
    pivot_comp = _construir_pivot(_preparar_df(dia_tipico_comparacion))

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=figsize, sharey=True)
    _dibujar_areas(ax2, pivot_comp, f"Fecha de referencia: {fecha_tipica_comparacion}")
    _dibujar_areas(ax1, pivot_est,  f"Fecha de referencia: {fecha_tipica}")

    handles, labels = ax1.get_legend_handles_labels()
    if not handles:
        handles, labels = ax2.get_legend_handles_labels()

    leg = fig.legend(handles, labels, title="Tecnología", loc="lower center",
                     ncol=len(labels), fontsize=8, title_fontsize=9,
                     frameon=True, bbox_to_anchor=(0.5, -0.05))
    _estilo_leyenda(leg)

    fig.subplots_adjust(left=0.06, right=0.98, top=0.84, bottom=0.20, wspace=0.08)
    _guardar_fig(fig, out_path, dpi=dpi)


# ══════════════════════════════════════════════════════════════════════════════
# 3) SPREAD CMG
# ══════════════════════════════════════════════════════════════════════════════

def graficar_spread_cmg(df_spread, out_path, figsize=(6.04, 4.29), dpi=300):
    columnas = {"nombre_cmg", "horas_solares", "horas_no_solares"}
    faltantes = columnas - set(df_spread.columns)
    if faltantes:
        raise ValueError(f"Faltan columnas: {sorted(faltantes)}")

    df_plot = df_spread.copy()
    if "spread_abs" in df_plot.columns:
        df_plot = df_plot.sort_values("spread_abs", ascending=False).reset_index(drop=True)

    x     = np.arange(len(df_plot))
    width = 0.38
    fig, ax = plt.subplots(figsize=figsize)

    ax.bar(x - width/2, df_plot["horas_solares"],    width=width,
           label="Horas solares",    color=MUTED["c2"], edgecolor="white", linewidth=0.8)
    ax.bar(x + width/2, df_plot["horas_no_solares"], width=width,
           label="Horas no solares", color=MUTED["c1"], edgecolor="white", linewidth=0.8)

    ax.set_title("CMG promedio: horas solares vs no solares", fontsize=12,
                 fontweight="bold", color=FONT_COLOR)
    ax.set_xlabel("Barra CMG", fontsize=10, color=FONT_COLOR)
    ax.set_ylabel("CMG promedio ($/kWh)", fontsize=10, color=FONT_COLOR)
    ax.set_xticks(x)
    ax.set_xticklabels(df_plot["nombre_cmg"], rotation=35, ha="right", fontsize=8)
    _estilo_ax(ax)

    leg = ax.legend(frameon=True, fontsize=8)
    _estilo_leyenda(leg)

    fig.subplots_adjust(left=0.09, right=0.98, top=0.87, bottom=0.28)
    _guardar_fig(fig, out_path, dpi=dpi)


# ══════════════════════════════════════════════════════════════════════════════
# 4) BOXPLOT VERTIMIENTOS + LÍNEA TOTAL
# ══════════════════════════════════════════════════════════════════════════════

def graficar_boxplot_vertimientos_con_total(df_all, out_path,
                                             figsize=(10.4, 4.8), dpi=300):
    df_box = df_all.copy()
    df_box["periodo"] = pd.to_datetime(df_box["periodo"], errors="coerce").dt.strftime("%Y-%m").astype(str)
    order_periodos = sorted(df_box["periodo"].dropna().unique().tolist())

    df_line = (df_box.groupby("periodo", as_index=False)["vertimiento"]
               .sum().rename(columns={"vertimiento": "kwh_total"}))
    df_line["periodo"] = pd.Categorical(df_line["periodo"], categories=order_periodos, ordered=True)
    df_line = df_line.sort_values("periodo")
    pos_map = {p: i for i, p in enumerate(order_periodos)}
    x_pos   = df_line["periodo"].astype(str).map(pos_map).astype(float)

    fig, ax = plt.subplots(figsize=figsize)
    ax2 = ax.twinx()

    ax2.plot(x_pos, df_line["kwh_total"].values,
             linewidth=2.2, color=MUTED["c2"], alpha=0.9, zorder=1, label="Total mensual")
    ax2.set_ylabel("kWh total mensual", color=FONT_COLOR)
    ax2.yaxis.set_major_formatter(FuncFormatter(_fmt_thousands))
    ax2.tick_params(colors=FONT_COLOR)
    ax2.grid(False)
    ax2.patch.set_alpha(0)
    sns.despine(ax=ax2, top=True, left=True)

    sns.boxplot(data=df_box, x="periodo", y="vertimiento", order=order_periodos,
                ax=ax, width=0.55, fliersize=2.2, linewidth=1.0,
                color=MUTED["c1"], zorder=3)

    ax.set_title("Vertimientos por empresa en cada mes (kWh)", fontsize=12,
                 fontweight="bold", color=FONT_COLOR)
    ax.set_xlabel("Periodo", fontsize=10, color=FONT_COLOR)
    ax.set_ylabel("kWh (distribución)", fontsize=10, color=FONT_COLOR)
    _estilo_ax(ax)
    ax.tick_params(axis="x", rotation=45)

    leg = ax2.legend(loc="upper right", bbox_to_anchor=(0.99, 0.99),
                     frameon=True, fontsize=9)
    _estilo_leyenda(leg)

    fig.subplots_adjust(left=0.06, right=0.96, top=0.88, bottom=0.20)
    _guardar_fig(fig, out_path, dpi=dpi)


# ══════════════════════════════════════════════════════════════════════════════
# 5) INYECTADA VS VERTIDA
# ══════════════════════════════════════════════════════════════════════════════

def graficar_inyectada_vertida(df_gx_ver_iny, out_path,
                                figsize=(5.19, 4.86), dpi=300):
    columnas = {"periodo", "inyeccion", "vertimiento"}
    faltantes = columnas - set(df_gx_ver_iny.columns)
    if faltantes:
        raise ValueError(f"Faltan columnas: {sorted(faltantes)}")

    df_plot = df_gx_ver_iny.copy().sort_values("periodo", ascending=True).reset_index(drop=True)
    x     = np.arange(len(df_plot))
    width = 0.55

    fig, ax = plt.subplots(figsize=figsize)

    ax.bar(x, df_plot["inyeccion"],  width=width, label="Inyecciones",
           color=MUTED["c2"], edgecolor="white", linewidth=0.8)
    ax.bar(x, df_plot["vertimiento"], width=width, bottom=df_plot["inyeccion"],
           label="Vertimientos", color=MUTED["c1"], edgecolor="white", linewidth=0.8)

    ax.set_title("Energía inyectada y vertida por periodo", fontsize=12,
                 fontweight="bold", color=FONT_COLOR)
    ax.set_xlabel("Periodo", fontsize=10, color=FONT_COLOR)
    ax.set_ylabel("Energía mWh", fontsize=10, color=FONT_COLOR)
    ax.set_xticks(x)
    ax.set_xticklabels(df_plot["periodo"], rotation=35, ha="right", fontsize=8)
    _estilo_ax(ax)

    leg = ax.legend(frameon=True, fontsize=8)
    _estilo_leyenda(leg)

    fig.subplots_adjust(left=0.09, right=0.98, top=0.87, bottom=0.28)
    _guardar_fig(fig, out_path, dpi=dpi)


# ══════════════════════════════════════════════════════════════════════════════
# 6) EVOLUCIÓN INYECCIÓN BESS
# ══════════════════════════════════════════════════════════════════════════════

def graficar_evolucion_inyeccion_bess(df_gx_real, df_gx_real_comparacion,
                                       out_path, figsize=(5.90, 4.14), dpi=300):
    COLORES = {0: MUTED["c2"], 1: MUTED["c1"]}

    df = evolucion_inyeccion_bess(df_gx_real, df_gx_real_comparacion)
    if df.empty:
        fig, ax = plt.subplots(figsize=figsize)
        ax.text(0.5, 0.5, "Sin datos de inyección BESS", ha="center", va="center",
                fontsize=11, color=FONT_COLOR)
        ax.axis("off")
        _guardar_fig(fig, out_path, dpi=dpi)
        return

    anios      = sorted(df["anio"].unique())
    trimestres = sorted(df["trimestre"].unique())
    x      = np.arange(len(trimestres))
    ancho  = 0.35
    offset = ancho / 2

    fig, ax = plt.subplots(figsize=figsize)

    for i, anio in enumerate(anios):
        df_a    = df[df["anio"] == anio].set_index("trimestre")
        valores = [df_a.loc[t, "inyeccion_retiro"] if t in df_a.index else 0 for t in trimestres]
        pos     = x - offset + i * ancho
        color   = COLORES.get(i, MUTED["c5"])

        bars = ax.bar(pos, valores, width=ancho, color=color, label=anio, linewidth=0)

        max_val = max([v for v in valores if v > 0], default=1)
        for bar, val in zip(bars, valores):
            if val > 0:
                ax.text(bar.get_x() + bar.get_width() / 2,
                        bar.get_height() + max_val * 0.01,
                        f"{val:,.0f}".replace(",", "."),
                        ha="center", va="bottom", fontsize=7, color=FONT_COLOR,
                        fontfamily=FONT_FAMILY)

    ax.set_xticks(x)
    ax.set_xticklabels([t.split("Q")[1] + "T" for t in trimestres], fontsize=8)
    ax.set_xlabel("Trimestre", fontsize=9, color=FONT_COLOR)
    ax.set_title("Evolución trimestral de inyección BESS", fontsize=11,
                 fontweight="bold", color=FONT_COLOR)
    ax.set_ylabel("Inyección (MWh)", fontsize=9, color=FONT_COLOR)
    _estilo_ax(ax)

    leg = ax.legend(title="Año", fontsize=8, title_fontsize=9,
                    frameon=True, loc="upper left")
    _estilo_leyenda(leg)

    fig.subplots_adjust(left=0.08, right=0.98, top=0.88, bottom=0.15)
    _guardar_fig(fig, out_path, dpi=dpi)


# ══════════════════════════════════════════════════════════════════════════════
# 7) EVOLUCIÓN VERTIMIENTO
# ══════════════════════════════════════════════════════════════════════════════

def graficar_evolucion_vertimiento(df_vertimientos, df_vertimientos_comparacion,
                                    out_path, figsize=(5.90, 3.85), dpi=300):
    COLORES = {0: MUTED["c2"], 1: MUTED["c1"]}

    df = evolucion_vertimiento(df_vertimientos, df_vertimientos_comparacion)
    if df.empty:
        fig, ax = plt.subplots(figsize=figsize)
        ax.text(0.5, 0.5, "Sin datos de vertimiento", ha="center", va="center",
                fontsize=11, color=FONT_COLOR)
        ax.axis("off")
        _guardar_fig(fig, out_path, dpi=dpi)
        return

    anios      = sorted(df["anio"].unique())
    trimestres = sorted(df["trimestre"].unique())
    x      = np.arange(len(trimestres))
    ancho  = 0.35
    offset = ancho / 2

    fig, ax = plt.subplots(figsize=figsize)

    for i, anio in enumerate(anios):
        df_a    = df[df["anio"] == anio].set_index("trimestre")
        valores = [df_a.loc[t, "vertimiento"] if t in df_a.index else 0 for t in trimestres]
        pos     = x - offset + i * ancho
        color   = COLORES.get(i, MUTED["c5"])

        bars = ax.bar(pos, valores, width=ancho, color=color, label=anio, linewidth=0)

        max_val = max([v for v in valores if v > 0], default=1)
        for bar, val in zip(bars, valores):
            if val > 0:
                ax.text(bar.get_x() + bar.get_width() / 2,
                        bar.get_height() + max_val * 0.01,
                        f"{val:,.0f}".replace(",", "."),
                        ha="center", va="bottom", fontsize=7, color=FONT_COLOR,
                        fontfamily=FONT_FAMILY)

    ax.set_xticks(x)
    ax.set_xticklabels([t.split("Q")[1] + "T" for t in trimestres], fontsize=8)
    ax.set_xlabel("Trimestre", fontsize=9, color=FONT_COLOR)
    ax.set_title("Evolución trimestral de vertimientos", fontsize=11,
                 fontweight="bold", color=FONT_COLOR)
    ax.set_ylabel("Vertimiento (MWh)", fontsize=9, color=FONT_COLOR)
    _estilo_ax(ax)

    leg = ax.legend(title="Año", fontsize=8, title_fontsize=9,
                    frameon=True, loc="upper left")
    _estilo_leyenda(leg)

    fig.subplots_adjust(left=0.08, right=0.98, top=0.88, bottom=0.15)
    _guardar_fig(fig, out_path, dpi=dpi)


# ══════════════════════════════════════════════════════════════════════════════
# ORQUESTADOR PRINCIPAL
# ══════════════════════════════════════════════════════════════════════════════

def generar_graficas(
    df_vertimientos,
    df_vertimientos_comparacion,
    df_spread,
    df_cmg,
    df_cmg_comparacion,
    df_dia_tipico,
    df_dia_tipico_comparacion,
    fecha_tipica,
    fecha_tipica_comparacion,
    df_gx_ver_iny,
    ppt_path,
    df_top_vertimiento,
    gx_real,
    gx_real_comparacion,
    outdir="outputs",
    dpi=300,
):
    from utils.insercion_graficos import insertar_graficos_ppt, get_figsize

    os.makedirs(outdir, exist_ok=True)
    _setup_theme()

    rename_map = {
        "CRUCERO_______220": "Barra crucero 200kV",
        "AJAHUEL_______500": "Barra Alto Jahuel 500kV",
        "P.MONTT_______220": "Barra Puerto Montt 220kV",
    }

    # ── 1) CMG con mapa ───────────────────────────────────────────────────────
    gdf_reg = generar_mapa_regiones(SHP_REGIONES)
    df_c    = df_cmg.copy()
    df_c["fecha_hora"] = pd.to_datetime(df_c["fecha_hora"], errors="coerce")
    df_c = df_c.sort_values("fecha_hora")
    df_c["nombre_cmg"] = df_c["nombre_cmg"].replace(rename_map)

    df_com = df_cmg_comparacion.copy()
    df_com["fecha_hora"] = pd.to_datetime(df_com["fecha_hora"], errors="coerce")
    df_com = df_com.sort_values("fecha_hora")
    df_com["nombre_cmg"] = df_com["nombre_cmg"].replace(rename_map)

    graficar_cmg_con_mapa(
        df_cmg=df_c, df_cmg_comparacion=df_com,
        gdf_reg=gdf_reg, bar_points=BAR_POINTS,
        out_path=os.path.join(outdir, "cmg.png"),
        figsize=get_figsize("img_cmg"), dpi=dpi,
    )

    # ── 2) Día típico ─────────────────────────────────────────────────────────
    graficar_gx_tipico(
        df_dia_tipico=df_dia_tipico,
        dia_tipico_comparacion=df_dia_tipico_comparacion,
        fecha_tipica=fecha_tipica,
        fecha_tipica_comparacion=fecha_tipica_comparacion,
        out_path=os.path.join(outdir, "gx_tipico.png"),
        figsize=get_figsize("img_dia_tipico"), dpi=dpi,
    )

    # ── 3) Spread CMG ─────────────────────────────────────────────────────────
    df_spread_plot = df_spread.copy()
    df_spread_plot["nombre_cmg"] = df_spread_plot["nombre_cmg"].replace(rename_map)
    graficar_spread_cmg(
        df_spread=df_spread_plot,
        out_path=os.path.join(outdir, "spread_cmg.png"),
        figsize=get_figsize("img_spread"), dpi=dpi,
    )

    # ── 4) Boxplot vertimientos ───────────────────────────────────────────────
    graficar_boxplot_vertimientos_con_total(
        df_all=df_vertimientos,
        out_path=os.path.join(outdir, "boxplot.png"),
        dpi=dpi,
    )

    # ── 5) Inyectada vs vertida ───────────────────────────────────────────────
    graficar_inyectada_vertida(
        df_gx_ver_iny=df_gx_ver_iny,
        out_path=os.path.join(outdir, "inyec_vert.png"),
        figsize=get_figsize("img_inyecciones_vertimientos"), dpi=dpi,
    )

    # ── 6) Evolución inyección BESS ───────────────────────────────────────────
    graficar_evolucion_inyeccion_bess(
        df_gx_real=gx_real,
        df_gx_real_comparacion=gx_real_comparacion,
        out_path=os.path.join(outdir, "inyecciones_bess.png"),
        figsize=get_figsize("img_inyeccion_bess"), dpi=dpi,
    )

    # ── 7) Evolución vertimientos ─────────────────────────────────────────────
    graficar_evolucion_vertimiento(
        df_vertimientos=df_vertimientos,
        df_vertimientos_comparacion=df_vertimientos_comparacion,
        out_path=os.path.join(outdir, "evolucion_vertimiento.png"),
        figsize=get_figsize("img_evolucion_vertimientos"), dpi=dpi,
    )

    # ── 8) Tabla top vertimientos ─────────────────────────────────────────────
    #Limpieza excepcional del df
    df_top_vertimiento["nombre_central"] = df_top_vertimiento["nombre_central"].str.replace("Pfv", "", regex=False)
    df_top_vertimiento["nombre_central"] = df_top_vertimiento["nombre_central"].str.replace("PFV", "", regex=False)
    print(df_top_vertimiento.head())


    insertar_top_vertimiento(ppt_path, 2, df_top_vertimiento=df_top_vertimiento)

    # ── 9) Insertar todas las imágenes en el PPT ──────────────────────────────
    insertar_graficos_ppt(ppt_path, outdir)
