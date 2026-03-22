from pptx import Presentation
from pptx.util import Pt
import pandas as pd
from pathlib import Path
import comtypes.client
import os
import time
import comtypes.client
import win32gui
import win32process
import psutil

# ==========================================================
# Helpers
# ========================================================== 
def set_textbox_text(slide, nombre, texto):
    for shape in slide.shapes:
        if shape.name == nombre:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = texto
                    return

def get_shape_by_name(slide, nombre):
    for shape in slide.shapes:
        if shape.name == nombre:
            return shape
    return None

def exportar_ppt_a_pdf(pptx_path, pdf_path):
    pptx_path = os.path.abspath(pptx_path)
    pdf_path  = os.path.abspath(pdf_path)

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    try:
        prs = powerpoint.Presentations.Open(pptx_path)
        prs.SaveAs(pdf_path, 32)  # 32 = ppSaveAsPDF
        prs.Close()
        print(f"📄 PDF exportado: {pdf_path}")
    finally:
        powerpoint.Quit()


# ==========================================================
# Funciones
# ========================================================== 

def insertar_top_vertimiento(ppt_path: Path,slide_idx, df_top_vertimiento=pd.DataFrame, top=10) -> pd.DataFrame:
    #Esta función abre el ppt e inserta las filas del df vertimiento en la tabla correspondiente del reporte
    # slide_idx: índice de la diapositiva
    # shape_idx_ indice del shape que contiene la tabla 

    # Importa la ppt y define el número de slide
    prs= Presentation(ppt_path)
    slide=prs.slides[slide_idx]

    #Obtiene el shape que contiene la tabla
    shape=get_shape_by_name(slide,"tabla_top")

    
    table=shape.table
    #Formateo numéricos
    for col in df_top_vertimiento.columns:
        if pd.api.types.is_numeric_dtype(df_top_vertimiento[col]):
            df_top_vertimiento[col] = df_top_vertimiento[col].map(
                lambda v: f"{v:,.0f}".replace(",",".") if 
                pd.notnull(v) else ""
            )    

    # Escribir fila por fila, columna por columna
    # row 0 generalmente es el encabezado — los datos van desde row 1
    for row_idx, row_data in enumerate(df_top_vertimiento.itertuples(index=False), start=1):
        if row_idx >= len(table.rows):
            break  # no exceder las filas existentes en la tabla
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            
            # Preservar tamaño de fuente del formato original
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)

    prs.save(ppt_path)
    print(f"Modifcado en {ppt_path}")


def insertar_periodo_estudio(ppt_path: Path, periodo_estudio: str):
    prs = Presentation(ppt_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name == "periodo_estudio":
                # Solo reemplaza el Run 1 — preserva formato de ambos runs
                para = shape.text_frame.paragraphs[0]
                if len(para.runs) >= 2:
                    para.runs[1].text = periodo_estudio
    prs.save(ppt_path)

def insertar_texto_con_placeholders(ppt_path: Path, kpis: dict) -> None:
    """
    Reemplaza placeholders tipo {{kpis.nombre}} preservando formato.
    Maneja placeholders partidos entre múltiples runs.
    """
    from pptx.util import Pt
    import copy

    prs = Presentation(ppt_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if not para.runs:
                    continue

                # Reconstruir texto completo del párrafo
                texto_completo = "".join(run.text for run in para.runs)

                # Verificar si hay algún placeholder en este párrafo
                tiene_placeholder = any(
                    f"{{{{kpis.{key}}}}}" in texto_completo
                    for key in kpis
                )
                if not tiene_placeholder:
                    continue

                # Reemplazar todos los placeholders
                for key, valor in kpis.items():
                    placeholder = f"{{{{kpis.{key}}}}}"
                    texto_completo = texto_completo.replace(placeholder, str(valor))

                # Preservar formato del primer run y escribir texto completo
                run0 = para.runs[0]
                fmt  = {
                    "bold":      run0.font.bold,
                    "italic":    run0.font.italic,
                    "size":      run0.font.size,
                    "color":     run0.font.color.rgb if run0.font.color and run0.font.color.type else None,
                    "underline": run0.font.underline,
                    "name":      run0.font.name,
                }

                run0.text = texto_completo

                # Reaplicar formato al run0
                run0.font.bold      = fmt["bold"]
                run0.font.italic    = fmt["italic"]
                run0.font.size      = fmt["size"]
                run0.font.underline = fmt["underline"]
                run0.font.name      = fmt["name"]
                if fmt["color"]:
                    run0.font.color.rgb = fmt["color"]

                # Limpiar runs restantes
                for run in para.runs[1:]:
                    run.text = ""

    prs.save(ppt_path)
    print(f"Texto insertado en {ppt_path}")