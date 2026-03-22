# utils/inserta_graficos_ppt.py
from pptx import Presentation
from pathlib import Path
from utils.config_loader import get_config

_cfg = get_config()

PLACEHOLDER_IMG_MAP = _cfg["ppt"]["imagenes"]

PLACEHOLDER_DIMS = {
    k: tuple(v) for k, v in _cfg["ppt"]["dimensiones"].items()
}

TEXTO_CONTEXTO_PREFIJOS = tuple(_cfg["ppt"]["texto_contexto_prefijos"])

# DPI para todos los gráficos
TARGET_DPI = _cfg["visualizacion"]["dpi"]
def get_figsize(placeholder_name: str, dpi: int = TARGET_DPI) -> tuple[float, float]:
    """
    Retorna (width_in, height_in) exactas del placeholder.
    Con este figsize + el dpi indicado, la imagen sale en los píxeles
    exactos del espacio en PPT — sin escalado.
    """
    dims = PLACEHOLDER_DIMS.get(placeholder_name)
    if dims is None:
        raise ValueError(f"'{placeholder_name}' no tiene dimensiones registradas.")
    return dims


def _buscar_shape_recursivo(shapes, nombre: str):
    """Busca un shape por nombre incluyendo dentro de grupos."""
    for shape in shapes:
        if shape.name == nombre:
            return shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP = 6
            encontrado = _buscar_shape_recursivo(shape.shapes, nombre)
            if encontrado:
                return encontrado
    return None



def insertar_graficos_ppt(ppt_path: Path, img_dir: Path, margen: float = None) -> None:
    if margen is None:
        margen = _cfg["ppt"]["margen_pulgadas"]
    """
    margen: margen interior en pulgadas por cada lado.
    Inserta las imágenes dentro de los placeholders y trae al frente
    todos los shapes de texto de contexto (títulos, ejes).
    """
    ppt_path   = Path(ppt_path)
    img_dir    = Path(img_dir)
    prs        = Presentation(ppt_path)
    margen_emu = int(margen * 914400)

    for slide_num, slide in enumerate(prs.slides, 1):

        # ── Insertar imágenes ─────────────────────────────────────
        for nombre, archivo in PLACEHOLDER_IMG_MAP.items():

            shape = _buscar_shape_recursivo(slide.shapes, nombre)
            if shape is None:
                continue

            img_file = img_dir / archivo
            if not img_file.exists():
                print(f"  ⚠️  Imagen no encontrada: {archivo} (slide {slide_num})")
                continue

            slide.shapes.add_picture(
                str(img_file),
                left=shape.left    + margen_emu,
                top=shape.top      + margen_emu,
                width=shape.width  - margen_emu * 2,
                height=shape.height - margen_emu * 2,
            )
            print(f"  ✅ {nombre} → {archivo} (slide {slide_num})")

        # ── Traer al frente títulos y etiquetas de ejes ───────────
        shapes_al_frente = [
            shape for shape in slide.shapes
            if any(shape.name.startswith(p) for p in TEXTO_CONTEXTO_PREFIJOS)
        ]
        for shape in shapes_al_frente:
            shape._element.getparent().append(shape._element)
            print(f"  ↑  '{shape.name}' traído al frente")

    prs.save(ppt_path)
    print(f"\nPPT guardado: {ppt_path}")