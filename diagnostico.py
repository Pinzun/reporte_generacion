from pptx import Presentation
path=r"C:\Users\Pablo\OneDrive\Documents\Trabajo\proyectos\reporte_generacion\data\raw\templates\template_reporte.pptx"
prs = Presentation(path)
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.name == "periodo_estudio":
            print(f"\n=== {shape.name} (slide) ===")
            for pi, para in enumerate(shape.text_frame.paragraphs):
                print(f"  Párrafo {pi}: '{para.text}'")
                for ri, run in enumerate(para.runs):
                    print(f"    Run {ri}: '{run.text}' | bold={run.font.bold} | size={run.font.size} | color={run.font.color.rgb if run.font.color.type else 'heredado'}")