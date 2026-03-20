from pptx import Presentation
from pathlib import Path

def explorar_ppt(ppt_path):
    # Aseguramos que ppt_path sea un objeto Path
    ppt_path = Path(ppt_path)
    prs = Presentation(ppt_path)
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*50}")
        print(f"SLIDE {slide_num}")
        print(f"{'='*50}")
        for shape in slide.shapes:
            # Convertir EMU (914400 = 1 inch) a pulgadas
            ancho = shape.width  / 914400
            alto  = shape.height / 914400
            print(f"  {shape.name:<30} | {ancho:.2f} x {alto:.2f} in")

if __name__ == "__main__":
    ppt_path = r"C:\Users\pinzunza\projectos_codigos\reporte_generacion\data\raw\templates\template_reporte.pptx"
    explorar_ppt(ppt_path)
